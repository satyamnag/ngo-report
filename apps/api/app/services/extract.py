"""Extract readable text from uploaded research documents.

Supported: txt, pdf, docx, xlsx, pptx. Images (png/jpg) return None — they are
stored for the user but contain no extractable text (vision description is a
future extension).
"""

import io


def extract_text(filename: str, data: bytes, content_type: str) -> str | None:
    ext = (filename or "").rsplit(".", 1)[-1].lower() if "." in (filename or "") else ""
    if ext == "txt":
        return data.decode("utf-8", errors="replace")[: 200_000]

    if ext == "pdf":
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        parts = []
        for page in reader.pages:
            parts.append(page.extract_text() or "")
        return "\n".join(parts)[: 200_000]

    if ext == "docx":
        from docx import Document

        doc = Document(io.BytesIO(data))
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text for c in row.cells))
        return "\n".join(parts)[: 200_000]

    if ext == "xlsx":
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"[Sheet: {ws.title}]")
            for row in ws.iter_rows(values_only=True):
                vals = [str(v) for v in row if v is not None]
                if vals:
                    parts.append(" | ".join(vals))
        return "\n".join(parts)[: 200_000]

    if ext == "pptx":
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts)[: 200_000]

    if ext in ("png", "jpg", "jpeg", "webp", "gif"):
        return None  # image: stored, no extractable text

    # Unknown type: try plain text if it looks text-like.
    if content_type and content_type.startswith("text/"):
        return data.decode("utf-8", errors="replace")[: 200_000]
    return None